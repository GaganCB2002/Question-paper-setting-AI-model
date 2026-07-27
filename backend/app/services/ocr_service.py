import io
import os
import re
import shutil
import aiofiles
from typing import Optional, Tuple
from pathlib import Path

import pytesseract
from PIL import Image, ImageEnhance, ImageFilter
import cv2
import numpy as np
from pdf2image import convert_from_path
from langdetect import detect, DetectorFactory, LangDetectException

from app.config import settings

DetectorFactory.seed = 0


class OcrService:
    def __init__(self):
        self.tesseract_cmd = settings.TESSERACT_CMD
        self.langs = settings.TESSERACT_LANGS
        self.dpi = settings.OCR_DPI
        if os.path.exists(self.tesseract_cmd) or shutil.which(self.tesseract_cmd):
            pytesseract.pytesseract.tesseract_cmd = self.tesseract_cmd

    async def process_image(self, image_bytes: bytes) -> dict:
        try:
            image = Image.open(io.BytesIO(image_bytes))
            return await self._ocr_image(image)
        except Exception as e:
            return {
                "raw_text": "",
                "cleaned_text": "",
                "language": "unknown",
                "confidence": 0,
                "page_count": 0,
                "tables": [],
                "headings": [],
                "words": [],
                "error": str(e),
            }

    async def process_pdf(self, pdf_path: str) -> dict:
        try:
            images = convert_from_path(pdf_path, dpi=self.dpi)
        except Exception as e:
            return {
                "raw_text": "",
                "cleaned_text": "",
                "language": "unknown",
                "confidence": 0,
                "page_count": 0,
                "tables": [],
                "headings": [],
                "error": f"PDF conversion failed: {str(e)}",
            }
        full_text = []
        all_tables = []
        all_headings = []
        total_confidence = 0.0
        page_data = []

        for page_num, image in enumerate(images, 1):
            result = await self._ocr_image(image, page_num=page_num)
            full_text.append(result.get("cleaned_text", ""))
            total_confidence += result.get("confidence", 0)
            if result.get("tables"):
                all_tables.extend(result["tables"])
            if result.get("headings"):
                all_headings.extend(result["headings"])
            page_data.append({
                "page": page_num,
                "text": result.get("cleaned_text", ""),
                "tables": result.get("tables", []),
                "headings": result.get("headings", []),
            })

        combined_text = "\n\n".join(full_text)
        lang = await self.detect_language(combined_text) if combined_text.strip() else "unknown"
        avg_confidence = total_confidence / len(images) if images else 0

        return {
            "raw_text": combined_text,
            "cleaned_text": await self._clean_text(combined_text),
            "language": lang,
            "confidence": round(avg_confidence, 2),
            "page_count": len(images),
            "tables": all_tables,
            "headings": all_headings,
            "pages": page_data,
        }

    async def _ocr_image(
        self, image: Image.Image, page_num: int = None
    ) -> dict:
        image = await self._enhance_image(image)
        image = await self._deskew(image)
        image = await self._remove_noise(image)

        custom_config = f"--oem 3 --psm 6 -l {self.langs}"

        data = pytesseract.image_to_data(
            image, config=custom_config, output_type=pytesseract.Output.DICT
        )
        text = pytesseract.image_to_string(image, config=custom_config)
        confidence_values = [
            int(conf) for conf in data["conf"] if conf != "-1"
        ]
        avg_confidence = (
            sum(confidence_values) / len(confidence_values)
            if confidence_values else 0
        )

        cleaned = await self._clean_text(text)
        tables = await self._extract_tables(image, data)
        headings = await self._extract_headings(data)

        words = data["text"]
        word_confidences = data["conf"]
        word_data = []
        for i, word in enumerate(words):
            if word.strip():
                word_data.append({
                    "word": word,
                    "confidence": int(word_confidences[i]) if word_confidences[i] != "-1" else 0,
                    "x": data["left"][i],
                    "y": data["top"][i],
                    "width": data["width"][i],
                    "height": data["height"][i],
                })

        return {
            "raw_text": text,
            "cleaned_text": cleaned,
            "confidence": round(float(avg_confidence), 2),
            "tables": tables,
            "headings": headings,
            "words": word_data,
            "page": page_num,
        }

    async def _enhance_image(self, image: Image.Image) -> Image.Image:
        image = image.convert("L")
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(2.0)
        enhancer = ImageEnhance.Sharpness(image)
        image = enhancer.enhance(2.0)
        return image

    async def _deskew(self, image: Image.Image) -> Image.Image:
        try:
            img_cv = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
            gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
            coords = np.column_stack(np.where(gray > 0))
            if len(coords) == 0:
                return image
            angle = cv2.minAreaRect(coords)[-1]
            if angle < -45:
                angle = 90 + angle
            angle = -angle
            if abs(angle) > 0.5:
                h, w = img_cv.shape[:2]
                center = (w // 2, h // 2)
                matrix = cv2.getRotationMatrix2D(center, angle, 1.0)
                rotated = cv2.warpAffine(
                    img_cv, matrix, (w, h),
                    flags=cv2.INTER_CUBIC,
                    borderMode=cv2.BORDER_REPLICATE
                )
                return Image.fromarray(cv2.cvtColor(rotated, cv2.COLOR_BGR2RGB))
            return image
        except Exception:
            return image

    async def _remove_noise(self, image: Image.Image) -> Image.Image:
        image = image.filter(ImageFilter.MedianFilter(size=3))
        img_cv = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
        img_cv = cv2.fastNlMeansDenoising(img_cv, None, 10, 7, 21)
        _, img_cv = cv2.threshold(img_cv, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        return Image.fromarray(cv2.cvtColor(img_cv, cv2.COLOR_BGR2RGB))

    async def _clean_text(self, text: str) -> str:
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'[^\w\s\u0C80-\u0CFF.,!?;:()\-"\'/\\%@#&*+={}\[\]|~`<>^$]', '', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = text.strip()
        return text

    async def _extract_tables(
        self, image: Image.Image, data: dict
    ) -> list[dict]:
        tables = []
        img_cv = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
        gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
        edges = cv2.Canny(gray, 50, 150, apertureSize=3)
        lines = cv2.HoughLinesP(
            edges, 1, np.pi / 180, 100,
            minLineLength=100, maxLineGap=10
        )
        if lines is not None:
            tables.append({
                "type": "detected",
                "line_count": len(lines),
                "confidence": "medium",
            })
        return tables

    async def _extract_headings(self, data: dict) -> list[dict]:
        headings = []
        for i, text in enumerate(data["text"]):
            if text.strip() and int(data["height"][i]) > 20:
                headings.append({
                    "text": text.strip(),
                    "font_size": data["height"][i],
                    "page": data.get("page_num", 1),
                })
        return headings

    async def detect_language(self, text: str) -> str:
        if not text.strip():
            return "unknown"
        try:
            lang = detect(text[:500])
            lang_map = {
                "kn": "kannada",
                "en": "english",
                "ml": "malayalam",
                "ta": "tamil",
                "te": "telugu",
                "hi": "hindi",
            }
            return lang_map.get(lang, lang)
        except LangDetectException:
            kannada_chars = len(re.findall(r'[\u0C80-\u0CFF]', text))
            english_chars = len(re.findall(r'[a-zA-Z]', text))
            if kannada_chars > english_chars:
                return "kannada"
            elif english_chars > kannada_chars:
                return "english"
            return "mixed"

    async def extract_text_from_docx(self, file_path: str) -> str:
        from docx import Document
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    async def extract_text_from_pptx(self, file_path: str) -> str:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)

    async def extract_text_from_xlsx(self, file_path: str) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True)
        texts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    texts.append(row_text)
        return "\n".join(texts)

    async def extract_text_from_txt(self, file_path: str) -> str:
        async with aiofiles.open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return await f.read()
